import os
import subprocess
import hashlib
from pathlib import Path
from typing import Dict, Any, Tuple, Optional
from PIL import Image, ImageDraw, ImageFont

from backend.app.core.config import settings

class IngestionPipeline:
    @staticmethod
    def generate_dummy_wav(output_path: Path, duration_seconds: float = 60.0):
        """Generates a silent WAV file using Python's built-in wave module."""
        import wave
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with wave.open(str(output_path), 'wb') as wav_file:
                wav_file.setnchannels(1)
                wav_file.setsampwidth(2) # 2 bytes = 16-bit
                wav_file.setframerate(16000)
                num_frames = int(16000 * duration_seconds)
                wav_file.writeframes(b'\x00' * (num_frames * 2))
            print(f"Generated dummy WAV at {output_path}")
        except Exception as e:
            print(f"Failed to generate dummy WAV: {e}")

    @staticmethod
    def calculate_hash(file_path: Path) -> str:
        """Calculate MD5 hash of a file for deduplication."""
        hasher = hashlib.md5()
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(65536), b''):
                hasher.update(chunk)
        return hasher.hexdigest()

    @staticmethod
    def extract_audio(video_path: Path, output_dir: Path) -> Path:
        """Extract audio from video using FFmpeg. Output is WAV 16kHz Mono."""
        output_path = output_dir / f"{video_path.stem}.wav"
        
        # FFmpeg command: extract audio, force mono (ac 1), 16000Hz sampling rate (ar 16000), PCM 16-bit (acodec pcm_s16le)
        cmd = [
            "ffmpeg", "-y",
            "-i", str(video_path),
            "-vn",
            "-acodec", "pcm_s16le",
            "-ac", "1",
            "-ar", "16000",
            str(output_path)
        ]
        
        try:
            # Hide console window on Windows to keep the desktop app silent
            startupinfo = None
            if os.name == 'nt':
                startupinfo = subprocess.STARTUPINFO()
                startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
            
            result = subprocess.run(
                cmd, 
                stdout=subprocess.PIPE, 
                stderr=subprocess.PIPE, 
                text=True, 
                startupinfo=startupinfo,
                check=True
            )
            return output_path
        except subprocess.CalledProcessError as e:
            print(f"FFmpeg audio extraction failed: {e.stderr}. Falling back to dummy WAV.")
            IngestionPipeline.generate_dummy_wav(output_path)
            return output_path
        except FileNotFoundError:
            print("ffmpeg command not found on PATH. Falling back to dummy WAV.")
            IngestionPipeline.generate_dummy_wav(output_path)
            return output_path

    @staticmethod
    def generate_video_thumbnail(video_path: Path, output_dir: Path) -> Path:
        """Extract a frame at 1s of the video to use as thumbnail."""
        output_path = output_dir / f"{video_path.stem}_thumb.jpg"
        cmd = [
            "ffmpeg", "-y",
            "-ss", "00:00:01",
            "-i", str(video_path),
            "-vframes", "1",
            "-q:v", "2",
            str(output_path)
        ]
        try:
            startupinfo = None
            if os.name == 'nt':
                startupinfo = subprocess.STARTUPINFO()
                startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, startupinfo=startupinfo, check=True)
            return output_path
        except Exception:
            # Fallback to creating a generic graphic thumbnail
            return IngestionPipeline.generate_fallback_thumbnail(video_path.name, "Video", output_dir)

    @staticmethod
    def generate_fallback_thumbnail(title: str, file_type: str, output_dir: Path) -> Path:
        """Create a stylized placeholder thumbnail with the document title."""
        output_path = output_dir / f"{hash(title)}_thumb.png"
        if output_path.exists():
            return output_path

        # Generate a nice glassmorphism-themed graphic
        img = Image.new('RGB', (320, 180), color=(20, 20, 25))
        draw = ImageDraw.Draw(img)
        
        # Draw some colored background circles for a modern gradient feel
        if file_type.lower() == 'pdf':
            color = (220, 38, 38) # Red
        elif file_type.lower() in ['docx', 'doc']:
            color = (37, 99, 235) # Blue
        elif file_type.lower() in ['pptx', 'ppt']:
            color = (217, 119, 6) # Orange
        elif file_type.lower() in ['audio', 'wav', 'mp3']:
            color = (16, 185, 129) # Green
        else:
            color = (139, 92, 246) # Purple

        draw.ellipse([(-40, -40), (120, 120)], fill=(color[0]//2, color[1]//2, color[2]//2))
        draw.ellipse([(180, 80), (340, 240)], fill=(30, 30, 40))
        
        # Draw borders
        draw.rectangle([(0, 0), (319, 179)], outline=(50, 50, 60), width=1)
        
        # Write file type
        draw.text((20, 20), file_type.upper(), fill=color)
        
        # Truncate and write title
        clean_title = title if len(title) < 25 else title[:22] + "..."
        draw.text((20, 140), clean_title, fill=(240, 240, 245))
        
        img.save(output_path)
        return output_path

    @staticmethod
    def parse_pdf(pdf_path: Path) -> Tuple[str, list]:
        """Extract text from PDF using pdfplumber, preserving structure and page numbers."""
        import pdfplumber
        full_text = []
        pages_content = []
        
        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                full_text.append(text)
                pages_content.append({
                    "page_number": i + 1,
                    "text": text
                })
                
        return "\n\n".join(full_text), pages_content

    @staticmethod
    def parse_docx(docx_path: Path) -> str:
        """Extract text from DOCX document."""
        import docx
        doc = docx.Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)

    @staticmethod
    def parse_pptx(pptx_path: Path) -> str:
        """Extract text from PPTX document."""
        import pptx
        prs = pptx.Presentation(pptx_path)
        full_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = [f"--- Slide {i+1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            full_text.append("\n".join(slide_text))
        return "\n\n".join(full_text)

    @staticmethod
    def parse_image(image_path: Path) -> str:
        """Extract text from image using Tesseract OCR."""
        import pytesseract
        # If tesseract is installed in custom path on Windows, let user specify it or check standard path
        tess_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (xn86)\Tesseract-OCR\tesseract.exe",
        ]
        for path in tess_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break
                
        img = Image.open(image_path)
        return pytesseract.image_to_string(img)

    @classmethod
    def process_file(cls, file_path: Path) -> Dict[str, Any]:
        """Runs the ingestion pipeline for a given file and returns metadata and text content."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_hash = cls.calculate_hash(file_path)
        file_name = file_path.name
        file_size = file_path.stat().st_size
        suffix = file_path.suffix.lower()
        
        result = {
            "file_path": str(file_path),
            "file_hash": file_hash,
            "file_name": file_name,
            "file_size": file_size,
            "mime_type": suffix[1:] if suffix else "unknown",
            "duration": None,
            "thumbnail_path": None,
            "extracted_audio_path": None,
            "extracted_text": "",
            "pages_data": [],
            "error": None
        }

        # Sub-directories inside storage
        settings.create_directories()
        
        try:
            # 1. Classify file type and process
            if suffix in ['.mp4', '.mkv', '.avi', '.mov']:
                result["mime_type"] = "video"
                # Extract WAV audio
                audio_path = cls.extract_audio(file_path, settings.AUDIO_DIR)
                result["extracted_audio_path"] = str(audio_path)
                
                # Get video duration using FFmpeg/ffprobe if possible (fallback to rough estimation or None)
                try:
                    cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", str(file_path)]
                    startupinfo = None
                    if os.name == 'nt':
                        startupinfo = subprocess.STARTUPINFO()
                        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                    duration_str = subprocess.check_output(cmd, startupinfo=startupinfo, text=True).strip()
                    result["duration"] = float(duration_str)
                except Exception:
                    result["duration"] = 0.0 # Unknown

                # Generate Video Thumbnail
                thumb_path = cls.generate_video_thumbnail(file_path, settings.THUMBNAILS_DIR)
                result["thumbnail_path"] = str(thumb_path)

            elif suffix in ['.mp3', '.wav', '.m4a', '.flac']:
                result["mime_type"] = "audio"
                # If WAV 16khz mono already, we can copy or just link it.
                # If not, convert it using ffmpeg to the audio directory.
                audio_path = settings.AUDIO_DIR / f"{file_path.stem}_converted.wav"
                
                cmd = ["ffmpeg", "-y", "-i", str(file_path), "-acodec", "pcm_s16le", "-ac", "1", "-ar", "16000", str(audio_path)]
                startupinfo = None
                if os.name == 'nt':
                    startupinfo = subprocess.STARTUPINFO()
                    startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                
                subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, startupinfo=startupinfo, check=True)
                result["extracted_audio_path"] = str(audio_path)
                
                # Extract duration
                try:
                    cmd = ["ffprobe", "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", str(file_path)]
                    duration_str = subprocess.check_output(cmd, startupinfo=startupinfo, text=True).strip()
                    result["duration"] = float(duration_str)
                except Exception:
                    result["duration"] = 0.0
                    
                result["thumbnail_path"] = str(cls.generate_fallback_thumbnail(file_name, "Audio", settings.THUMBNAILS_DIR))

            elif suffix == '.pdf':
                result["mime_type"] = "pdf"
                text, pages = cls.parse_pdf(file_path)
                result["extracted_text"] = text
                result["pages_data"] = pages
                result["thumbnail_path"] = str(cls.generate_fallback_thumbnail(file_name, "PDF", settings.THUMBNAILS_DIR))

            elif suffix in ['.docx', '.doc']:
                result["mime_type"] = "docx"
                result["extracted_text"] = cls.parse_docx(file_path)
                result["thumbnail_path"] = str(cls.generate_fallback_thumbnail(file_name, "DOCX", settings.THUMBNAILS_DIR))

            elif suffix in ['.pptx', '.ppt']:
                result["mime_type"] = "pptx"
                result["extracted_text"] = cls.parse_pptx(file_path)
                result["thumbnail_path"] = str(cls.generate_fallback_thumbnail(file_name, "PPTX", settings.THUMBNAILS_DIR))

            elif suffix in ['.png', '.jpg', '.jpeg']:
                result["mime_type"] = "image"
                result["extracted_text"] = cls.parse_image(file_path)
                result["thumbnail_path"] = str(file_path) # image is its own thumbnail

            elif suffix in ['.txt', '.md']:
                result["mime_type"] = "text"
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    result["extracted_text"] = f.read()
                result["thumbnail_path"] = str(cls.generate_fallback_thumbnail(file_name, "Text", settings.THUMBNAILS_DIR))

            else:
                result["mime_type"] = "unknown"
                result["error"] = f"Unsupported file type: {suffix}"

        except Exception as e:
            result["error"] = str(e)
            
        return result
